"""
Script to generate PowerPoint presentation for CM2013 Sleep Scoring Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def add_section_header(prs, title):
    """Add a section header slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    slide.shapes.title.text = title
    return slide

def add_content_slide(prs, title, content_items):
    """Add a slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item['text']
        p.level = item.get('level', 0)
        if 'bullet' in item:
            p.text = item['bullet'] + " " + p.text

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Define table position and size
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(4)

    # Add table
    table = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height).table

    # Set headers
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 112, 192)
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.size = Pt(12)

    # Set rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(11)

    return slide

def add_flowchart_slide(prs, title, steps, description=None):
    """Add a slide with a vertical flowchart"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Starting position
    left = Inches(2.5)
    top = Inches(2)
    width = Inches(5)
    height = Inches(0.8)

    colors = [
        RGBColor(68, 114, 196),   # Blue
        RGBColor(112, 173, 71),   # Green
        RGBColor(237, 125, 49),   # Orange
        RGBColor(165, 165, 165),  # Gray
        RGBColor(255, 192, 0),    # Yellow
    ]

    for i, step in enumerate(steps):
        # Add shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top + i * (height + Inches(0.3)),
            width, height
        )

        # Style the shape
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[i % len(colors)]
        shape.line.color.rgb = RGBColor(0, 0, 0)
        shape.line.width = Pt(1)

        # Add text
        text_frame = shape.text_frame
        text_frame.text = step
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Add arrow (except for last step)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                left + width/2 - Inches(0.15),
                top + (i+1) * (height + Inches(0.3)) - Inches(0.25),
                Inches(0.3), Inches(0.2)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.color.rgb = RGBColor(100, 100, 100)

    if description:
        # Add description at bottom
        textbox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        text_frame = textbox.text_frame
        text_frame.text = description
        text_frame.paragraphs[0].font.size = Pt(12)
        text_frame.paragraphs[0].font.italic = True

    return slide

def add_comparison_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a slide comparing two approaches side by side"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Left box
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.8),
        Inches(4.25), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = RGBColor(252, 229, 205)  # Light orange
    left_box.line.color.rgb = RGBColor(237, 125, 49)
    left_box.line.width = Pt(2)

    # Right box
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.25), Inches(1.8),
        Inches(4.25), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(217, 234, 211)  # Light green
    right_box.line.color.rgb = RGBColor(112, 173, 71)
    right_box.line.width = Pt(2)

    # Left title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.25), Inches(0.5))
    tf = left_title_box.text_frame
    tf.text = left_title
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Right title
    right_title_box = slide.shapes.add_textbox(Inches(5.25), Inches(2), Inches(4.25), Inches(0.5))
    tf = right_title_box.text_frame
    tf.text = right_title
    tf.paragraphs[0].font.size = Pt(20)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Left content
    left_content_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(3.85), Inches(3.3))
    tf = left_content_box.text_frame
    for item in left_items:
        p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(12)
        p.space_after = Pt(8)

    # Right content
    right_content_box = slide.shapes.add_textbox(Inches(5.45), Inches(2.7), Inches(3.85), Inches(3.3))
    tf = right_content_box.text_frame
    for item in right_items:
        p = tf.add_paragraph()
        p.text = "✓ " + item
        p.font.size = Pt(12)
        p.space_after = Pt(8)

    return slide

def add_timeline_slide(prs, title, milestones):
    """Add a slide with a horizontal timeline"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Timeline line
    line = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(1), Inches(4),
        Inches(9), Inches(4)
    )
    line.line.color.rgb = RGBColor(68, 114, 196)
    line.line.width = Pt(3)

    # Add milestones
    num_milestones = len(milestones)
    spacing = Inches(8) / (num_milestones - 1) if num_milestones > 1 else 0

    colors = [
        RGBColor(68, 114, 196),
        RGBColor(112, 173, 71),
        RGBColor(237, 125, 49),
        RGBColor(165, 165, 165),
    ]

    for i, milestone in enumerate(milestones):
        x = Inches(1) + i * spacing

        # Milestone circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x - Inches(0.25), Inches(4) - Inches(0.25),
            Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i % len(colors)]
        circle.line.color.rgb = RGBColor(255, 255, 255)
        circle.line.width = Pt(2)

        # Milestone label (above)
        label_box = slide.shapes.add_textbox(x - Inches(0.8), Inches(2.5), Inches(1.6), Inches(0.6))
        tf = label_box.text_frame
        tf.text = milestone['label']
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Milestone date (below)
        date_box = slide.shapes.add_textbox(x - Inches(0.8), Inches(4.6), Inches(1.6), Inches(0.5))
        tf = date_box.text_frame
        tf.text = milestone['date']
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)

    return slide

def add_architecture_diagram(prs, title):
    """Add a system architecture diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    # Layer definitions
    layers = [
        {"name": "Data Layer", "color": RGBColor(68, 114, 196), "items": ["EDF Files", "XML Files", "Cache"]},
        {"name": "Processing Layer", "color": RGBColor(112, 173, 71), "items": ["Data Loading", "Preprocessing", "Features", "Selection"]},
        {"name": "Intelligence Layer", "color": RGBColor(237, 125, 49), "items": ["Classifiers", "Evaluation"]},
        {"name": "Output Layer", "color": RGBColor(165, 165, 165), "items": ["Results", "Visualization", "Reports"]},
    ]

    top = Inches(1.8)
    height = Inches(1)

    for i, layer in enumerate(layers):
        # Main box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), top + i * (height + Inches(0.2)),
            Inches(8), height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = layer["color"]
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.line.width = Pt(1)

        # Layer name
        name_box = slide.shapes.add_textbox(
            Inches(1.2), top + i * (height + Inches(0.2)) + Inches(0.1),
            Inches(7.6), Inches(0.3)
        )
        tf = name_box.text_frame
        tf.text = layer["name"]
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Items
        items_text = "  |  ".join(layer["items"])
        items_box = slide.shapes.add_textbox(
            Inches(1.2), top + i * (height + Inches(0.2)) + Inches(0.5),
            Inches(7.6), Inches(0.4)
        )
        tf = items_box.text_frame
        tf.text = items_text
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Arrow to next layer
        if i < len(layers) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.85), top + (i+1) * (height + Inches(0.2)) - Inches(0.15),
                Inches(0.3), Inches(0.15)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(100, 100, 100)
            arrow.line.color.rgb = RGBColor(100, 100, 100)

    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
                   "CM2013: Sleep Scoring Project",
                   "Biomedical Signal Processing\nAutomated Sleep Stage Classification")

    # Slide 2: Project Overview
    add_content_slide(prs, "Project Overview", [
        {'text': 'Goal: Develop an automatic sleep scoring system', 'level': 0},
        {'text': 'Multi-signal biosignal processing (EEG, EOG, EMG)', 'level': 1},
        {'text': 'Machine learning classification of 5 sleep stages', 'level': 1},
        {'text': 'Duration: 10 weeks with 4 iterations', 'level': 0},
        {'text': 'Team-based development (3 members per group)', 'level': 0},
        {'text': 'Deliverables:', 'level': 0},
        {'text': 'Working code (Python/MATLAB)', 'level': 1},
        {'text': 'Technical report (15 pages max)', 'level': 1},
        {'text': 'Project management documentation', 'level': 1},
    ])

    # Slide 3: Learning Objectives
    add_content_slide(prs, "Learning Objectives", [
        {'text': 'Apply signal processing techniques to real biomedical data', 'level': 0},
        {'text': 'Filtering, artifact removal, feature extraction', 'level': 1},
        {'text': 'Implement machine learning classifiers for pattern recognition', 'level': 0},
        {'text': 'k-NN, SVM, Random Forest', 'level': 1},
        {'text': 'Practice agile software development in teams', 'level': 0},
        {'text': 'Sprints, iterative development, task management', 'level': 1},
        {'text': 'Develop professional documentation skills', 'level': 0},
        {'text': 'Code documentation, technical reports, presentations', 'level': 1},
    ])

    # Slide 4: Assessment Criteria
    add_content_slide(prs, "Assessment Criteria", [
        {'text': 'Methodology and Code Quality - 50%', 'level': 0},
        {'text': 'Modular design, correct pipeline, documentation, testing', 'level': 1},
        {'text': 'Team Collaboration - 30%', 'level': 0},
        {'text': 'Regular updates, integration, ClickUp usage', 'level': 1},
        {'text': 'Report & Documentation - 20%', 'level': 0},
        {'text': 'Clear technical writing, comprehensive analysis', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Note: No fixed accuracy target - focus on process and learning!', 'level': 0},
    ])

    # Slide 5: Section Header - Development Methodology
    add_section_header(prs, "Development Methodology:\nAgile vs. Waterfall")

    # Slide 6: Waterfall vs Agile Comparison (VISUAL)
    add_comparison_slide(prs, "Waterfall vs. Agile Methodology",
        "Waterfall ❌",
        [
            "Sequential phases",
            "No feedback until end",
            "High risk (all-or-nothing)",
            "Rigid - difficult to adapt",
            "Late discovery of issues",
            "Clear structure upfront"
        ],
        "Agile ✓ (Our Approach)",
        [
            "Iterative cycles (sprints)",
            "Feedback after each iteration",
            "Low risk (incremental)",
            "Flexible - adapt to results",
            "Early problem detection",
            "Perfect for research projects"
        ]
    )

    # Slide 7: Sprint Lifecycle (VISUAL FLOWCHART)
    add_flowchart_slide(prs, "Sprint Lifecycle (2.5 weeks each)",
        [
            "Day 1-2: Sprint Planning",
            "Day 3-15: Development",
            "Day 16-17: Testing & Integration",
            "Day 18: Sprint Review"
        ],
        "Repeat for each of 4 iterations"
    )

    # Slide 8: Section Header - Iteration Planning
    add_section_header(prs, "Iteration Planning:\n4 Sprints Over 10 Weeks")

    # Slide 9: Visual Timeline
    add_timeline_slide(prs, "Project Timeline - Key Milestones",
        [
            {"label": "Iteration 1\nBasic EEG", "date": "Oct 31, 2025"},
            {"label": "Iteration 2\nEEG+EOG", "date": "Nov 19, 2025"},
            {"label": "Iteration 3\nEEG+EOG+EMG", "date": "Dec 5, 2025"},
            {"label": "Iteration 4\nFinal System", "date": "Dec 18, 2025"}
        ]
    )

    # Slide 10: Iteration Timeline Table
    add_table_slide(prs, "Iteration Timeline & Milestones",
                   ["Iteration", "Deadline", "Signals", "Features", "Classifier"],
                   [
                       ["1", "Oct 31, 2025", "EEG", "Time (16)", "k-NN"],
                       ["2", "Nov 19, 2025", "EEG+EOG", "Time+Freq (31)", "SVM"],
                       ["3", "Dec 5, 2025", "EEG+EOG+EMG", "Selected (30)", "Random Forest"],
                       ["4", "Dec 18, 2025", "All signals", "Optimized", "RF-optimized"]
                   ])

    # Slide 11: Iteration 1 Details
    add_content_slide(prs, "Iteration 1: Basic Pipeline (Oct 31)", [
        {'text': 'Focus: Get something working end-to-end', 'level': 0},
        {'text': 'Signals: EEG only (single or dual channel)', 'level': 0},
        {'text': 'Features: 16 time-domain features per channel', 'level': 0},
        {'text': 'Mean, median, std, variance, RMS, Hjorth parameters', 'level': 1},
        {'text': 'Classifier: k-Nearest Neighbors (k-NN)', 'level': 0},
        {'text': 'Simple, interpretable baseline', 'level': 1},
        {'text': 'Deliverables:', 'level': 0},
        {'text': 'Complete data loading pipeline', 'level': 1},
        {'text': 'Basic preprocessing (filtering, epoching)', 'level': 1},
        {'text': 'Working end-to-end classification', 'level': 1},
    ])

    # Slide 12: Iteration 2 Details
    add_content_slide(prs, "Iteration 2: Enhanced Processing (Nov 19)", [
        {'text': 'Focus: Add EOG and frequency-domain features', 'level': 0},
        {'text': 'Signals: EEG + EOG (eye movement detection)', 'level': 0},
        {'text': 'Features: ~31 features (time + frequency domain)', 'level': 0},
        {'text': 'Add frequency features: band powers, spectral entropy', 'level': 1},
        {'text': 'EOG-specific: eye movement characteristics', 'level': 1},
        {'text': 'Classifier: Support Vector Machine (SVM)', 'level': 0},
        {'text': 'Better handling of high-dimensional data', 'level': 1},
        {'text': 'Deliverables:', 'level': 0},
        {'text': 'Multi-signal processing capability', 'level': 1},
        {'text': 'Frequency-domain feature extraction', 'level': 1},
    ])

    # Slide 13: Iteration 3 Details
    add_content_slide(prs, "Iteration 3: Multi-Signal (Dec 5)", [
        {'text': 'Focus: Add EMG and implement feature selection', 'level': 0},
        {'text': 'Signals: EEG + EOG + EMG (muscle tone)', 'level': 0},
        {'text': 'Features: ~30 selected features (down from 50+)', 'level': 0},
        {'text': 'Feature selection: statistical tests, mutual information', 'level': 1},
        {'text': 'EMG features: muscle activity indicators', 'level': 1},
        {'text': 'Classifier: Random Forest', 'level': 0},
        {'text': 'Handles non-linear relationships, provides feature importance', 'level': 1},
        {'text': 'Deliverables:', 'level': 0},
        {'text': 'Complete multi-signal processing', 'level': 1},
        {'text': 'Intelligent feature selection', 'level': 1},
    ])

    # Slide 14: Iteration 4 Details
    add_content_slide(prs, "Iteration 4: Full System (Dec 18)", [
        {'text': 'Focus: Optimization and finalization', 'level': 0},
        {'text': 'Signals: All available channels optimally combined', 'level': 0},
        {'text': 'Features: Optimized feature set', 'level': 0},
        {'text': 'Cross-validation for robustness', 'level': 1},
        {'text': 'Classifier: Optimized Random Forest', 'level': 0},
        {'text': 'Hyperparameter tuning, ensemble methods', 'level': 1},
        {'text': 'Deliverables:', 'level': 0},
        {'text': 'Final competition submission', 'level': 1},
        {'text': 'Complete technical report', 'level': 1},
        {'text': 'Project documentation and presentation', 'level': 1},
    ])

    # Slide 15: Section Header - ClickUp
    add_section_header(prs, "Project Management:\nUsing ClickUp")

    # Slide 16: Why Use ClickUp?
    add_content_slide(prs, "Why Project Management Tools?", [
        {'text': 'Essential for team collaboration and coordination', 'level': 0},
        {'text': 'Track who is doing what and when', 'level': 1},
        {'text': 'Visualize progress and identify blockers', 'level': 1},
        {'text': 'Maintain accountability and transparency', 'level': 1},
        {'text': 'Professional skill development:', 'level': 0},
        {'text': 'Industry standard practice', 'level': 1},
        {'text': 'Critical for remote/distributed teams', 'level': 1},
        {'text': 'Required for assessment (grading checkpoints)', 'level': 0},
        {'text': 'Instructor reviews your ClickUp at each milestone', 'level': 1},
    ])

    # Slide 17: ClickUp Setup
    add_content_slide(prs, "ClickUp Setup (Project Manager Task)", [
        {'text': '1. Designate one team member as Project Manager', 'level': 0},
        {'text': 'Role can rotate between iterations', 'level': 1},
        {'text': '2. Create workspace: CM2013_Sleep_Scoring_Group[X]', 'level': 0},
        {'text': '3. Add all team members with edit access', 'level': 0},
        {'text': '4. ⚠️ MANDATORY: Add instructor as viewer', 'level': 0},
        {'text': '5. Create sprint folders:', 'level': 0},
        {'text': 'Iteration 1: Basic EEG (Due: Oct 31, 2025)', 'level': 1},
        {'text': 'Iteration 2: EEG+EOG (Due: Nov 19, 2025)', 'level': 1},
        {'text': 'Iteration 3: EEG+EOG+EMG (Due: Dec 5, 2025)', 'level': 1},
        {'text': 'Iteration 4: Full System (Due: Dec 18, 2025)', 'level': 1},
    ])

    # Slide 18: Task Organization
    add_content_slide(prs, "Task Organization in ClickUp", [
        {'text': 'Create tags for organization (free version):', 'level': 0},
        {'text': 'Priority: 🔴 HIGH, 🟡 MEDIUM, 🟢 LOW', 'level': 1},
        {'text': 'Signals: #EEG, #EOG, #EMG', 'level': 1},
        {'text': 'Components: #preprocessing, #features, #classification', 'level': 1},
        {'text': 'Status: #BLOCKED, #NEEDS-REVIEW, #BUG', 'level': 1},
        {'text': 'Task workflow:', 'level': 0},
        {'text': 'To Do → In Progress → Review → Testing → Complete', 'level': 1},
        {'text': 'Each task must have:', 'level': 0},
        {'text': 'Clear title: [Component] Specific action', 'level': 1},
        {'text': 'Assignee, due date, priority, description', 'level': 1},
    ])

    # Slide 19: Daily Standups
    add_content_slide(prs, "Daily Standups (via ClickUp)", [
        {'text': 'Each team member posts daily update as task comment:', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Format:', 'level': 0},
        {'text': '"Today: [what I did]', 'level': 1},
        {'text': 'Tomorrow: [what I\'ll do]', 'level': 1},
        {'text': 'Blockers: [any issues]"', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Benefits:', 'level': 0},
        {'text': 'Keep everyone informed without meetings', 'level': 1},
        {'text': 'Identify problems early', 'level': 1},
        {'text': 'Build accountability and momentum', 'level': 1},
        {'text': 'Tag PM if blocked: @mention', 'level': 1},
    ])

    # Slide 20: Grading Checkpoints
    add_content_slide(prs, "ClickUp Grading Checkpoints", [
        {'text': 'Instructor will review your ClickUp at:', 'level': 0},
        {'text': '', 'level': 0},
        {'text': '✓ October 31, 2025 - Iteration 1 complete', 'level': 0},
        {'text': '✓ November 19, 2025 - Iteration 2 complete', 'level': 0},
        {'text': '✓ December 5, 2025 - Iteration 3 complete', 'level': 0},
        {'text': '✓ December 18, 2025 - Final delivery', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'What is evaluated:', 'level': 0},
        {'text': 'Task organization and clarity', 'level': 1},
        {'text': 'Regular updates and progress', 'level': 1},
        {'text': 'Team communication and collaboration', 'level': 1},
        {'text': 'Problem-solving and adaptability', 'level': 1},
    ])

    # Slide 21: Section Header - Python Jumpstart
    add_section_header(prs, "Python Jumpstart:\nProject Structure & Usage")

    # Slide 22: What's Provided
    add_content_slide(prs, "Python Jumpstart: What's Provided", [
        {'text': '⚠️ Structure and examples ONLY - not complete solution!', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Provided:', 'level': 0},
        {'text': 'Modular project structure (src/ directory)', 'level': 1},
        {'text': 'Configuration system (config.py)', 'level': 1},
        {'text': 'Basic filter example (lowpass at 40Hz)', 'level': 1},
        {'text': '3 simple features (mean, median, std)', 'level': 1},
        {'text': 'Basic k-NN classifier with train/test split', 'level': 1},
        {'text': 'Caching system for efficiency', 'level': 1},
        {'text': 'Testing framework (pytest)', 'level': 1},
        {'text': 'Google Colab notebook', 'level': 1},
    ])

    # Slide 23: Project Structure
    add_content_slide(prs, "Python Project Structure", [
        {'text': 'Python/', 'level': 0},
        {'text': 'src/ - Core modules:', 'level': 1},
        {'text': 'data_loader.py - Load EDF/XML files', 'level': 2},
        {'text': 'preprocessing.py - Signal filtering', 'level': 2},
        {'text': 'feature_extraction.py - Extract features', 'level': 2},
        {'text': 'feature_selection.py - Select best features', 'level': 2},
        {'text': 'classification.py - ML classifiers', 'level': 2},
        {'text': 'visualization.py - Plot results', 'level': 2},
        {'text': 'report.py - Generate reports', 'level': 2},
        {'text': 'main.py - Training pipeline orchestration', 'level': 1},
        {'text': 'run_inference.py - Generate predictions', 'level': 1},
        {'text': 'config.py - Central configuration', 'level': 1},
        {'text': 'colab_notebook.ipynb - Run in Google Colab', 'level': 1},
    ])

    # Slide 24: System Architecture Diagram (VISUAL)
    add_architecture_diagram(prs, "System Architecture - Data Flow")

    # Slide 25: Configuration System
    add_content_slide(prs, "Configuration System (config.py)", [
        {'text': 'Central control for the entire pipeline', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Key settings:', 'level': 0},
        {'text': 'CURRENT_ITERATION (1-4) - Controls which features/algorithms', 'level': 1},
        {'text': 'USE_CACHE (True/False) - Speed up development', 'level': 1},
        {'text': 'File paths - Data directories', 'level': 1},
        {'text': 'Preprocessing parameters - Filter frequencies', 'level': 1},
        {'text': 'Model hyperparameters - Classifier settings', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Benefits:', 'level': 0},
        {'text': 'Easy to switch between iterations', 'level': 1},
        {'text': 'Consistent settings across team', 'level': 1},
    ])

    # Slide 26: Running the Pipeline
    add_content_slide(prs, "Running the Training Pipeline", [
        {'text': '1. Setup (first time only):', 'level': 0},
        {'text': 'pip install -r requirements.txt', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '2. Verify setup:', 'level': 0},
        {'text': 'python -m pytest tests/ -v', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '3. Run training pipeline:', 'level': 0},
        {'text': 'python main.py', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '4. Run inference (generate submission):', 'level': 0},
        {'text': 'python run_inference.py', 'level': 1},
        {'text': 'Creates submission.csv in data/ directory', 'level': 1},
    ])

    # Slide 27: Caching System
    add_content_slide(prs, "Caching System for Efficiency", [
        {'text': 'Why caching?', 'level': 0},
        {'text': 'Preprocessing is slow (minutes per file)', 'level': 1},
        {'text': 'Feature extraction takes time', 'level': 1},
        {'text': 'Don\'t repeat work during development', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'How it works:', 'level': 0},
        {'text': 'Results saved to cache/ directory', 'level': 1},
        {'text': 'Automatically reused on next run', 'level': 1},
        {'text': 'Control with USE_CACHE in config.py', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'When to clear cache:', 'level': 0},
        {'text': 'Changed preprocessing parameters', 'level': 1},
        {'text': 'Modified feature extraction code', 'level': 1},
    ])

    # Slide 28: Google Colab Notebook
    add_content_slide(prs, "Google Colab Notebook", [
        {'text': 'Alternative to local development', 'level': 0},
        {'text': 'File: colab_notebook.ipynb', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Features:', 'level': 0},
        {'text': 'No local setup required', 'level': 1},
        {'text': 'Free GPU access (if needed)', 'level': 1},
        {'text': 'Load from GitHub or Google Drive', 'level': 1},
        {'text': 'Run complete pipeline in browser', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Usage:', 'level': 0},
        {'text': 'Upload to Google Colab', 'level': 1},
        {'text': 'Follow cell-by-cell instructions', 'level': 1},
        {'text': 'View results inline', 'level': 1},
    ])

    # Slide 29: What Students Must Implement
    add_content_slide(prs, "What You Must Implement", [
        {'text': 'Iteration 1:', 'level': 0},
        {'text': 'Real EDF/XML file parsing', 'level': 1},
        {'text': 'Bandpass filtering (0.5-40 Hz)', 'level': 1},
        {'text': '13+ additional time-domain features', 'level': 1},
        {'text': 'Iteration 2:', 'level': 0},
        {'text': 'Multi-channel processing (EEG+EOG)', 'level': 1},
        {'text': 'Frequency-domain features (band powers)', 'level': 1},
        {'text': 'SVM hyperparameter tuning', 'level': 1},
        {'text': 'Iteration 3:', 'level': 0},
        {'text': 'EMG signal processing', 'level': 1},
        {'text': 'Feature selection algorithms', 'level': 1},
        {'text': 'Iteration 4:', 'level': 0},
        {'text': 'Cross-validation, optimization, final report', 'level': 1},
    ])

    # Slide 30: Section Header - Report
    add_section_header(prs, "Final Report:\nStructure & Requirements")

    # Slide 31: Report Structure
    add_content_slide(prs, "Technical Report (15 pages max)", [
        {'text': '1. Introduction (1 page)', 'level': 0},
        {'text': 'Problem statement, objectives', 'level': 1},
        {'text': '2. Methods (3 pages)', 'level': 0},
        {'text': 'Signal processing, features, classification', 'level': 1},
        {'text': '3. Results (4-5 pages)', 'level': 0},
        {'text': 'Performance metrics, confusion matrices, analysis', 'level': 1},
        {'text': '4. Discussion (3-4 pages)', 'level': 0},
        {'text': 'Interpretation, challenges, improvements', 'level': 1},
        {'text': '5. Conclusion (1 page)', 'level': 0},
        {'text': 'Summary, future work', 'level': 1},
        {'text': '6. References (1 page)', 'level': 0},
    ])

    # Slide 32: Report Content Guidelines
    add_content_slide(prs, "Report Content Focus", [
        {'text': 'Focus on the process, not just results', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Methods section should explain:', 'level': 0},
        {'text': 'WHY you chose specific approaches', 'level': 1},
        {'text': 'HOW you implemented them', 'level': 1},
        {'text': 'What alternatives you considered', 'level': 1},
        {'text': 'Results section should show:', 'level': 0},
        {'text': 'Progression through iterations', 'level': 1},
        {'text': 'Impact of different decisions', 'level': 1},
        {'text': 'Statistical analysis of performance', 'level': 1},
        {'text': 'Discussion should reflect:', 'level': 0},
        {'text': 'Critical thinking about your approach', 'level': 1},
        {'text': 'Learning from what didn\'t work', 'level': 1},
    ])

    # Slide 33: Section Header - Data
    add_section_header(prs, "Data & File Formats")

    # Slide 34: Data Structure
    add_content_slide(prs, "Data Organization", [
        {'text': 'data/training/ - EDF + XML files with labels', 'level': 0},
        {'text': 'Use for training and validation', 'level': 1},
        {'text': 'Both EDF and XML required for each recording', 'level': 1},
        {'text': 'data/holdout/ - EDF files only (no labels)', 'level': 0},
        {'text': 'Use for final predictions/competition', 'level': 1},
        {'text': 'data/sample/ - Small test dataset', 'level': 0},
        {'text': 'Quick testing during development', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Key signals in EDF files:', 'level': 0},
        {'text': 'EEG: C3-A2, C4-A1 (125 Hz, hardware high-pass 0.15 Hz)', 'level': 1},
        {'text': 'EOG: Left/Right (50 Hz, hardware high-pass 0.15 Hz)', 'level': 1},
        {'text': 'EMG: (125 Hz, hardware high-pass 0.15 Hz)', 'level': 1},
        {'text': 'ECG: (125 Hz), Respiration: Thor/Abdo (10 Hz), SpO2/HR (1 Hz)', 'level': 1},
    ])

    # Slide 35: Signal Details Table
    add_table_slide(prs, "Available Signals - Detailed Specifications",
                   ["Signal", "EDF Label", "Sample Rate", "Hardware Filter"],
                   [
                       ["EEG C3-A2", "EEG (sec)", "125 Hz", "HP 0.15 Hz"],
                       ["EEG C4-A1", "EEG", "125 Hz", "HP 0.15 Hz"],
                       ["EOG Left", "EOG(L)", "50 Hz", "HP 0.15 Hz"],
                       ["EOG Right", "EOG(R)", "50 Hz", "HP 0.15 Hz"],
                       ["EMG", "EMG", "125 Hz", "HP 0.15 Hz"],
                       ["ECG", "ECG", "125 Hz", "HP 0.15 Hz"],
                       ["Thorax RES", "Thor RES", "10 Hz", "HP 0.05 Hz"],
                       ["Abdomen RES", "Abdo RES", "10 Hz", "HP 0.05 Hz"],
                       ["SpO2", "SaO2", "1 Hz", "-"],
                       ["Heart Rate", "H.R.", "1 Hz", "-"]
                   ])

    # Slide 36: EDF and XML Formats
    add_content_slide(prs, "EDF and XML File Formats", [
        {'text': 'EDF (European Data Format):', 'level': 0},
        {'text': 'Standard for biosignals (EEG, EOG, EMG)', 'level': 1},
        {'text': 'Multiple channels with metadata', 'level': 1},
        {'text': 'Python: Use MNE library (mne.io.read_raw_edf)', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'XML (Compumedics Annotation Format):', 'level': 0},
        {'text': 'Sleep stage labels for each 30-second epoch', 'level': 1},
        {'text': 'Stages: Wake, N1, N2, N3, REM', 'level': 1},
        {'text': 'Python: Use xml.etree.ElementTree', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '📚 Reference: github.com/nsrr/edf-editor-translator/wiki', 'level': 0},
    ])

    # Slide 37: Signal Processing Considerations
    add_content_slide(prs, "Signal Processing Considerations", [
        {'text': 'Hardware Filtering Already Applied:', 'level': 0},
        {'text': 'EEG/EOG/EMG/ECG: High-pass 0.15 Hz', 'level': 1},
        {'text': 'Respiration: High-pass 0.05 Hz', 'level': 1},
        {'text': 'Design additional filters accordingly', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Different Sampling Rates:', 'level': 0},
        {'text': 'Primary signals (EEG/EMG/ECG): 125 Hz', 'level': 1},
        {'text': 'EOG signals: 50 Hz', 'level': 1},
        {'text': 'Respiration/Airflow: 10 Hz', 'level': 1},
        {'text': 'SpO2/Heart Rate: 1 Hz', 'level': 1},
        {'text': '', 'level': 0},
        {'text': '30-second epochs = different sample counts per signal', 'level': 0},
    ])

    # Slide 38: Section Header - Success Tips
    add_section_header(prs, "Tips for Success")

    # Slide 39: General Tips
    add_content_slide(prs, "Tips for Success", [
        {'text': 'Start early and work consistently', 'level': 0},
        {'text': 'Don\'t wait until deadlines', 'level': 1},
        {'text': 'Code first, optimize later', 'level': 0},
        {'text': 'Get something working, then improve', 'level': 1},
        {'text': 'Test often - run pipeline after changes', 'level': 0},
        {'text': 'Use caching to speed up development', 'level': 0},
        {'text': 'Document as you go - not at the end', 'level': 0},
        {'text': 'Communicate proactively with your team', 'level': 0},
        {'text': 'Over-communication is better than under-communication', 'level': 1},
        {'text': 'Ask for help early when stuck', 'level': 0},
        {'text': 'Celebrate wins and learn from failures', 'level': 0},
    ])

    # Slide 40: Team Roles
    add_content_slide(prs, "Team Organization (3 members)", [
        {'text': 'Project Manager:', 'level': 0},
        {'text': 'Coordination, ClickUp, integration, documentation', 'level': 1},
        {'text': 'Preprocessing Lead:', 'level': 0},
        {'text': 'Signal cleaning, filtering, artifact removal', 'level': 1},
        {'text': 'Feature Engineer:', 'level': 0},
        {'text': 'Feature extraction, selection, analysis', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Note: With 3 members, roles overlap!', 'level': 0},
        {'text': 'Everyone should contribute to multiple areas', 'level': 1},
        {'text': 'Cross-train and help each other', 'level': 1},
        {'text': 'ML/classification can be shared responsibility', 'level': 1},
    ])

    # Slide 41: Resources
    add_content_slide(prs, "Resources & Documentation", [
        {'text': 'Project documentation:', 'level': 0},
        {'text': 'PROJECT_GUIDE.md - Complete project guide', 'level': 1},
        {'text': 'CLAUDE.md - Codebase overview', 'level': 1},
        {'text': 'Python/README.md - Python jumpstart guide', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Key libraries:', 'level': 0},
        {'text': 'MNE - EEG/biosignal processing', 'level': 1},
        {'text': 'scikit-learn - Machine learning', 'level': 1},
        {'text': 'NumPy/SciPy - Signal processing', 'level': 1},
        {'text': '', 'level': 0},
        {'text': 'Support:', 'level': 0},
        {'text': 'Office hours, course forum, team members', 'level': 1},
    ])

    # Slide 42: Questions
    add_section_header(prs, "Questions?")

    # Slide 43: Summary
    add_content_slide(prs, "Summary - Key Takeaways", [
        {'text': '✓ Iterative agile development over 10 weeks', 'level': 0},
        {'text': '✓ 4 sprints with clear milestones and deadlines', 'level': 0},
        {'text': '✓ Assessment based on process, not just accuracy', 'level': 0},
        {'text': '✓ Use ClickUp for project management and collaboration', 'level': 0},
        {'text': '✓ Python jumpstart provides structure, you implement algorithms', 'level': 0},
        {'text': '✓ Focus on learning and continuous improvement', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Start with Iteration 1: Get basic pipeline working!', 'level': 0},
        {'text': '', 'level': 0},
        {'text': 'Good luck! 🚀', 'level': 0},
    ])

    return prs

if __name__ == "__main__":
    prs = create_presentation()
    prs.save("CM2013_Sleep_Scoring_Project_Lecture.pptx")
    print("✓ Presentation created: CM2013_Sleep_Scoring_Project_Lecture.pptx")
    print("  - 43 slides with visual enhancements:")
    print("    • Project overview and objectives")
    print("    • Agile vs. Waterfall methodology (COMPARISON DIAGRAM)")
    print("    • Sprint lifecycle (FLOWCHART)")
    print("    • Iteration planning and milestones (TIMELINE)")
    print("    • System architecture (LAYERED DIAGRAM)")
    print("    • ClickUp setup and usage")
    print("    • Python jumpstart structure")
    print("    • Report requirements")
    print("    • Detailed signal specifications (TABLE)")
    print("    • Data formats and resources")